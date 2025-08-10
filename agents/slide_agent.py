"""
Clean slides/flashcards generation with template-based approach.
"""

import json
import asyncio
import re
import os
from typing import Dict, Any, Optional, List
from pathlib import Path

try:
    import google.generativeai as genai
    GENAI_AVAILABLE = True
except ImportError:
    GENAI_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from core.config import settings

class SlidesAgent:
    """Generate clean, professional flashcards using template-based approach."""
    
    def __init__(self):
        self.api_key = settings.GEMINI_API_KEY
        self.model = None
        self.is_initialized = False
        self.templates = {
            "frustrated": {"bg_color": (173, 216, 230), "accent": (70, 130, 180)},    # Light blue, calming
            "confused": {"bg_color": (255, 248, 220), "accent": (255, 140, 0)},       # Light yellow, clarity
            "happy": {"bg_color": (240, 248, 255), "accent": (30, 144, 255)},        # Alice blue, energetic
            "neutral": {"bg_color": (248, 249, 250), "accent": (52, 73, 94)}          # Light gray, professional
        }
    
    async def initialize(self):
        """Initialize Gemini 2.5 Pro model."""
        if self.is_initialized:
            return
            
        if not GENAI_AVAILABLE:
            print("✗ Google Generative AI not available")
            return
            
        if not self.api_key:
            print("✗ Gemini API key not found in config")
            return
            
        try:
            genai.configure(api_key=self.api_key)
            self.model = genai.GenerativeModel('gemini-2.5-pro')
            self.is_initialized = True
            print("✓ Gemini 2.5 Pro initialized for clean flashcards")
        except Exception as e:
            print(f"✗ Gemini initialization failed: {e}")
    
    async def generate_flashcards(
        self,
        content_chunks: List[str],
        emotion_context: Optional[str] = None,
        cards_per_chunk: int = 1
    ) -> Dict[str, Any]:
        """Generate clean, professional flashcards from content."""
        if not self.is_initialized:
            await self.initialize()
        
        if not self.model:
            return self._mock_flashcards(content_chunks, emotion_context)
        
        try:
            prompt = self._build_clean_prompt(content_chunks, emotion_context, cards_per_chunk)
            response = await asyncio.to_thread(
                self.model.generate_content, prompt
            )
            
            flashcards_data = self._parse_response(response.text)
            await self._save_flashcards(flashcards_data)
            
            if PPTX_AVAILABLE:
                ppt_path = await self._create_professional_pptx(flashcards_data)
                flashcards_data["ppt_file_path"] = str(ppt_path)
            
            return flashcards_data
            
        except Exception as e:
            print(f"✗ Flashcards generation failed: {e}")
            return self._mock_flashcards(content_chunks, emotion_context)
    
    def _build_clean_prompt(self, content_chunks: List[str], emotion_context: str, cards_per_chunk: int) -> str:
        """Build focused prompt for clean flashcard generation."""
        
        content_text = "\n\n".join([f"Chunk {i+1}: {chunk}" for i, chunk in enumerate(content_chunks)])
        total_cards = len(content_chunks) * cards_per_chunk
        
        return f"""
<thinking>
I need to create {total_cards} high-quality educational flashcards for a {emotion_context or 'neutral'} learner.

Key requirements:
- One clear concept per card
- Questions test understanding, not memorization  
- Clean, professional format
- Emotion-appropriate difficulty level
- Educational effectiveness over complexity
</thinking>

<reflect>
Are these questions clear and specific?
Do they test real understanding?
Is the difficulty appropriate for someone who is {emotion_context or 'neutral'}?
Will these actually help someone learn the material?
</reflect>

Create exactly {total_cards} educational flashcards from the content below. Focus on clarity, educational value, and appropriate difficulty for a {emotion_context or 'neutral'} learner.

Content:
{content_text}

Requirements:
- One concept per card
- Clear, specific questions
- Complete but concise answers
- Educational hints when helpful
- Professional, clean format

Output as valid JSON within <output></output> tags:

<output>
{{
  "metadata": {{
    "title": "Educational Flashcards",
    "total_cards": {total_cards},
    "emotion_context": "{emotion_context or 'neutral'}",
    "difficulty": "appropriate"
  }},
  "cards": [
    {{
      "id": 1,
      "question": "Clear, specific question",
      "answer": "Complete, concise answer",
      "hint": "Helpful hint if needed",
      "topic": "Main concept covered",
      "difficulty": "easy|medium|hard"
    }}
  ]
}}
</output>
"""
    
    def _parse_response(self, response_text: str) -> Dict[str, Any]:
        """Parse and validate flashcards JSON."""
        try:
            pattern = r"<output>(.*?)</output>"
            match = re.search(pattern, response_text, re.DOTALL)
            
            if match:
                json_str = match.group(1).strip()
                data = json.loads(json_str)
                
                if self._validate_structure(data):
                    return data
                    
        except (json.JSONDecodeError, Exception) as e:
            print(f"✗ Parsing failed: {e}")
        
        return self._mock_flashcards(["parsing failed"])
    
    def _validate_structure(self, data: Dict[str, Any]) -> bool:
        """Validate flashcards structure."""
        if "metadata" not in data or "cards" not in data:
            return False
            
        cards = data.get("cards", [])
        if not cards:
            return False
            
        for card in cards:
            required_keys = ["id", "question", "answer", "topic"]
            if not all(key in card for key in required_keys):
                return False
                
        return True
    
    async def _save_flashcards(self, data: Dict[str, Any]) -> None:
        """Save flashcards JSON."""
        try:
            output_dir = Path("outputs")
            output_dir.mkdir(exist_ok=True)
            
            file_path = output_dir / "flashcards.json"
            with open(file_path, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
                
            print(f"✓ Flashcards saved to {file_path}")
            
        except Exception as e:
            print(f"✗ Save failed: {e}")
    
    async def _create_professional_pptx(self, data: Dict[str, Any]) -> Optional[Path]:
        """Create clean, professional PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            print("✗ python-pptx not available")
            return None
            
        try:
            prs = Presentation()
            
            emotion = data.get("metadata", {}).get("emotion_context", "neutral")
            theme = self.templates.get(emotion, self.templates["neutral"])
            
            self._create_title_slide(prs, data, theme)
            
            cards = data.get("cards", [])
            for card in cards:
                self._create_flashcard_slide(prs, card, theme)
            
            output_dir = Path("outputs")
            output_dir.mkdir(exist_ok=True)
            ppt_path = output_dir / "flashcards.pptx"
            
            prs.save(str(ppt_path))
            print(f"✓ Professional PowerPoint saved to {ppt_path}")
            
            return ppt_path
            
        except Exception as e:
            print(f"✗ PowerPoint generation failed: {e}")
            return None
    
    def _create_title_slide(self, prs: Presentation, data: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Create clean title slide."""
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        
        bg_color = theme["bg_color"]
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = data.get("metadata", {}).get("title", "Educational Flashcards")
        subtitle_shape.text = f"Total Cards: {data.get('metadata', {}).get('total_cards', 0)} | Emotion Context: {data.get('metadata', {}).get('emotion_context', 'neutral').title()}"
        
        self._format_title_text(title_shape, theme)
        self._format_subtitle_text(subtitle_shape, theme)
    
    def _create_flashcard_slide(self, prs: Presentation, card: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Create individual flashcard slide with clean layout."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        bg_color = theme["bg_color"]
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
        
        self._add_card_header(slide, card, theme)
        self._add_question_section(slide, card, theme)
        self._add_answer_section(slide, card, theme)
        self._add_hint_section(slide, card, theme)
        self._add_footer(slide, card, theme)
    
    def _add_card_header(self, slide, card: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Add clean header with card info."""
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        header_frame = header_box.text_frame
        
        p = header_frame.paragraphs[0]
        p.text = f"Card {card.get('id', 1)} - {card.get('topic', 'Concept')}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*theme["accent"])
    
    def _add_question_section(self, slide, card: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Add question with clean formatting."""
        q_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(2))
        q_frame = q_box.text_frame
        q_frame.margin_left = Inches(0.2)
        q_frame.margin_right = Inches(0.2)
        q_frame.margin_top = Inches(0.1)
        q_frame.word_wrap = True
        
        p = q_frame.paragraphs[0]
        p.text = f"Q: {card.get('question', '')}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(33, 37, 41)
        p.alignment = PP_ALIGN.LEFT
    
    def _add_answer_section(self, slide, card: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Add answer with proper spacing."""
        a_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2.5))
        a_frame = a_box.text_frame
        a_frame.margin_left = Inches(0.2)
        a_frame.margin_right = Inches(0.2)
        a_frame.margin_top = Inches(0.1)
        a_frame.word_wrap = True
        
        p = a_frame.paragraphs[0]
        p.text = f"A: {card.get('answer', '')}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(52, 58, 64)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
    
    def _add_hint_section(self, slide, card: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Add hint if available."""
        if not card.get('hint'):
            return
            
        hint_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(1))
        hint_frame = hint_box.text_frame
        hint_frame.margin_left = Inches(0.2)
        hint_frame.word_wrap = True
        
        p = hint_frame.paragraphs[0]
        p.text = f"💡 Hint: {card.get('hint', '')}"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(108, 117, 125)
        p.alignment = PP_ALIGN.LEFT
    
    def _add_footer(self, slide, card: Dict[str, Any], theme: Dict[str, Any]) -> None:
        """Add difficulty indicator."""
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.2), Inches(1.5), Inches(0.5))
        footer_frame = footer_box.text_frame
        
        p = footer_frame.paragraphs[0]
        p.text = card.get('difficulty', 'medium').upper()
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        
        difficulty_colors = {"easy": (40, 167, 69), "medium": (255, 193, 7), "hard": (220, 53, 69)}
        color = difficulty_colors.get(card.get('difficulty', 'medium'), difficulty_colors["medium"])
        p.font.color.rgb = RGBColor(*color)
    
    def _format_title_text(self, shape, theme: Dict[str, Any]) -> None:
        """Format title with professional styling."""
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*theme["accent"])
    
    def _format_subtitle_text(self, shape, theme: Dict[str, Any]) -> None:
        """Format subtitle with clean styling."""
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(108, 117, 125)
    
    def _mock_flashcards(self, content_chunks: List[str], emotion_context: str = "neutral") -> Dict[str, Any]:
        """Generate clean mock flashcards."""
        return {
            "metadata": {
                "title": "Neural Networks & AI Concepts",
                "total_cards": min(len(content_chunks), 5),
                "emotion_context": emotion_context,
                "difficulty": "appropriate"
            },
            "cards": [
                {
                    "id": 1,
                    "question": "What is the core function of a neural network?",
                    "answer": "To process data through interconnected layers of nodes that learn patterns by adjusting connection weights based on training examples.",
                    "hint": "Think about how neurons in the brain connect and strengthen their connections through learning.",
                    "topic": "Neural Network Fundamentals",
                    "difficulty": "easy"
                },
                {
                    "id": 2,
                    "question": "How do transformers process sequences differently from RNNs?",
                    "answer": "Transformers use self-attention to process all elements simultaneously, while RNNs process elements sequentially one at a time.",
                    "hint": "Consider parallel vs sequential processing approaches.",
                    "topic": "Transformer Architecture", 
                    "difficulty": "medium"
                }
            ]
        }

slides_agent = SlidesAgent()

async def test_slides_agent():
    """Test clean slides generation."""
    print("Testing Clean Slides Agent")
    print("=" * 40)
    
    test_chunks = [
        "Neural networks consist of interconnected layers of nodes that process information by adjusting connection weights through training.",
        "Transformers use self-attention mechanisms to process entire sequences simultaneously rather than sequentially like RNNs.",
        "Backpropagation calculates gradients to update network weights, minimizing the difference between predicted and actual outputs.",
        "Convolutional Neural Networks use specialized filters to detect patterns in grid-like data such as images.",
        "Transfer learning allows pre-trained models to be adapted for new tasks with minimal additional training data."
    ]
    
    emotions = ["neutral", "frustrated", "happy", "confused"]
    
    for emotion in emotions:
        print(f"\nTesting: {emotion} learner")
        print("-" * 25)
        
        try:
            result = await slides_agent.generate_flashcards(
                content_chunks=test_chunks[:3],  # Limit for clean testing
                emotion_context=emotion,
                cards_per_chunk=1
            )
            
            metadata = result.get("metadata", {})
            cards = result.get("cards", [])
            
            print(f"✓ Generated: {metadata.get('title', 'Unknown')}")
            print(f"  Cards: {len(cards)}")
            print(f"  Context: {metadata.get('emotion_context', 'Unknown')}")
            
            if result.get("ppt_file_path"):
                print(f"  PowerPoint: {result['ppt_file_path']}")
                
            if cards:
                sample = cards[0]
                print(f"  Sample Q: {sample.get('question', '')[:50]}...")
                
        except Exception as e:
            print(f"✗ Test failed: {e}")
    
    print("\n" + "=" * 40)
    print("✓ Clean slides testing complete")
    print("Check outputs/flashcards.json and outputs/flashcards.pptx")

if __name__ == "__main__":
    asyncio.run(test_slides_agent())
